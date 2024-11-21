# Analyze and adjust based on the existing slide content
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.util import Inches


def get_shapes(template_slide):
    # If the slide is found, place the image at the exact position of an existing placeholder
    if template_slide:
        # Find a picture placeholder or an example shape
        image_positions = []
        text_shapes = []
        width = None
        height = None
        page_shape = None
        for shape in template_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if not image_positions:
                    width = shape.width
                    height = shape.height
                    # print(height)
                image_positions.append((shape.left, shape.top))

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_shapes.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # 確定這是頁碼框 (頁腳的一部分)
                if shape.placeholder_format.type == 10:  # 頁碼類型
                    # 複製頁碼框的位置和大小
                    page_shape = shape.left, shape.top, shape.width, shape.height
            else:
                # Handle other shape types or skip them
                print(f"Skipping unsupported shape type: {shape.shape_type}")
        return image_positions, width, height, text_shapes
    else:
        # If no specific placeholder, use default position and size
        image_positions = [
            (Inches(1), Inches(1)),  # Left-top
            (Inches(5), Inches(1)),  # Right-top
            (Inches(1), Inches(4)),  # Left-bottom
            (Inches(5), Inches(4))  # Right-bottom
        ]
        width, height = Inches(3), Inches(2)  # Image size
        print("no template!")


def add_new_slide(presentation, template_slide, image_dir_path, image_name, new_slide_title):
    # Add a new slide based on the template
    slide_layout = presentation.slide_layouts[5]  # Use blank layout
    new_slide = presentation.slides.add_slide(slide_layout)

    image_positions, width, height, text_shapes = get_shapes(template_slide)

    new_shape = new_slide.shapes.add_textbox(Inches(12.85), Inches(7.05), Inches(1), Inches(0.5))
    page_number = len(presentation.slides)  # 新頁碼
    new_shape.text = str(page_number)  # 新頁碼內容
    # print(new_shape.text)

    if image_positions and width and height and text_shapes:

        # Insert images into the slide
        for image_type, position in zip(image_types, image_positions):
            image_path = f'{image_dir_path}/{image_type}/{image_name}'
            left, top = position
            try:
                new_slide.shapes.add_picture(image_path, left, top, width, height)
                # print(image_path)
            except FileNotFoundError:
                print(f"Image not found: {image_path}")

        for text_shape in text_shapes:
            # Copy text box
            left, top, width, height = text_shape.left, text_shape.top, text_shape.width, text_shape.height
            new_shape = new_slide.shapes.add_textbox(left, top, width, height)
            new_shape.text = text_shape.text_frame.text
            # print((new_shape.text))
            new_shape.text_frame.paragraphs[0].runs[0].font.size = text_shape.text_frame.paragraphs[0].runs[0].font.size

        title_shape = new_slide.shapes.title  # 尋找標題框
        if title_shape:
            title_shape.text = new_slide_title  # 設定標題

        # Save the modified presentation
        presentation.save(output_pptx_path)
        print(output_pptx_path)


if __name__ == "__main__":
    # File paths
    pptx_path = r'C:/Users/Alice/OneDrive/lab/fusionProject/progress1120_GT_虹惠.pptx'
    output_pptx_path = 'C:/Users/Alice/OneDrive/lab/fusionProject/updated_progress1120_GT_虹惠.pptx'
    # Path to the image extracted from the zip
    image_types = ["gray_i1", "gray_i2", "GT", "fusion"]
    template_page = 3

    # Reload the PowerPoint presentation
    presentation = Presentation(pptx_path)
    # Look for a slide with a specific placeholder or text indicating "ep8"
    template_slide = presentation.slides[template_page - 1]
    # 啟用頁碼
    presentation.core_properties.slide_number = 1  # 頁碼從 1 開始

    image_dir_path1 = ['D:/GT_1118/seblock_ep8']
    new_slide_title1 = "SeBlock Model - Epoch 8"
    image_name1 = 'image_14.png'
    for image_dir_path in image_dir_path1:
        add_new_slide(presentation, template_slide, image_dir_path, image_name1, new_slide_title1)

    image_dir_path2 = ['D:/GT_1118/multi_scale_ep1', 'D:/GT_1118/multi_scale_ep4', 'D:/GT_1118/multi_scale_ep8']
    new_slide_title2 = [f"Multi Scale - Epoch {i}" for i in [1, 4, 8]]
    image_name2 = ['image_2.png', 'image_4.png', 'image_2.png']
    for image_dir_path, image_name, new_slide_title in zip(image_dir_path2, image_name2, new_slide_title2):
        add_new_slide(presentation, template_slide, image_dir_path, image_name, new_slide_title)

    image_dir_path3 = ['D:/GT_1118/perceptual_guide_ep1', 'D:/GT_1118/perceptual_guide_ep4',
                       'D:/GT_1118/perceptual_guide_ep8']
    new_slide_title3 = [f"Perceptual Guide - Epoch {i}" for i in [1, 4, 8]]
    image_name3 = ['image_0.png', 'image_13.png', 'image_9.png']
    for image_dir_path, image_name, new_slide_title in zip(image_dir_path3, image_name3, new_slide_title3):
        add_new_slide(presentation, template_slide, image_dir_path, image_name, new_slide_title)

